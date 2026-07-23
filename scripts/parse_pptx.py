# parse_pptx.py - 从 PPTX 提取文本与备注
# 用法: python parse_pptx.py <pptx_path> <output_md_path>
from pptx import Presentation
import sys, os

def extract_pptx(pptx_path, output_path):
    prs = Presentation(pptx_path)
    lines = ["# 绪论\n"]
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            lines.append(f"## 第 {i} 页\n")
            for t in texts:
                lines.append(t)
                lines.append("")
        # 备注
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"> **备注**：{notes}\n")
    content = "\n".join(lines)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"✅ 已提取 {len(prs.slides)} 页 -> {output_path}")

if __name__ == "__main__":
    pptx_path = sys.argv[1]
    output_path = sys.argv[2]
    extract_pptx(pptx_path, output_path)
